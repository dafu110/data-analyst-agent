from __future__ import annotations

from io import BytesIO

from backend.pdf_exporter import markdown_to_pdf


def markdown_to_pptx(markdown: str) -> bytes:
    try:
        from pptx import Presentation
    except ImportError as exc:  # pragma: no cover
        raise RuntimeError("PPTX 导出需要安装 python-pptx：pip install -e .[prod]") from exc

    presentation = Presentation()
    title_slide_layout = presentation.slide_layouts[0]
    bullet_layout = presentation.slide_layouts[1]
    lines = [line for line in markdown.splitlines() if line.strip()]
    title = next((line.lstrip("# ").strip() for line in lines if line.startswith("#")), "数据分析报告")

    slide = presentation.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = title
    slide.placeholders[1].text = "Data Analyst Agent"

    current_title = "关键内容"
    bullets: list[str] = []
    for line in lines:
        if line.startswith("## "):
            if bullets:
                add_bullet_slide(presentation, bullet_layout, current_title, bullets)
                bullets = []
            current_title = line[3:].strip()
        elif line.startswith("- "):
            bullets.append(line[2:].strip())
            if len(bullets) >= 7:
                add_bullet_slide(presentation, bullet_layout, current_title, bullets)
                bullets = []

    if bullets:
        add_bullet_slide(presentation, bullet_layout, current_title, bullets)

    buffer = BytesIO()
    presentation.save(buffer)
    return buffer.getvalue()


def add_bullet_slide(presentation, layout, title: str, bullets: list[str]) -> None:
    slide = presentation.slides.add_slide(layout)
    slide.shapes.title.text = title
    body = slide.placeholders[1].text_frame
    body.clear()
    for index, bullet in enumerate(bullets):
        paragraph = body.paragraphs[0] if index == 0 else body.add_paragraph()
        paragraph.text = bullet[:180]
        paragraph.level = 0
